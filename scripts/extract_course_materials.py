"""Extract all PS369 course materials (.docx, .pptx, .pdf) to plain text.

Idempotent: re-running this script picks up any newly added files and overwrites
existing extracts. Run this whenever you add new material under CourseMaterials/.

Output structure mirrors the input:
  CourseMaterials/PS369_Handouts/plain_text/<stem>.txt
  CourseMaterials/PS369_LectureMaterials/plain_text/<stem>.txt
  CourseMaterials/PS369_StudyGuides/plain_text/<stem>.txt
  CourseMaterials/syllabus_plain_text.txt

Readings (PS369_Readings/) are full academic texts and not currently extracted;
cite them by title/author/page rather than extracting.
"""
from __future__ import annotations

import pathlib
import subprocess
import sys

ROOT = pathlib.Path('/Users/jordangans-morse/Dropbox/PostSovietApp/CourseMaterials')

try:
    import docx  # python-docx
    from pptx import Presentation
except ImportError as e:
    sys.exit(f'Missing dependency: {e}. Install with: pip3 install python-docx python-pptx')


def extract_docx(path: pathlib.Path) -> str:
    d = docx.Document(str(path))
    lines: list[str] = []
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            lines.append(t)
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            lines.append(' | '.join(cells))
    return '\n'.join(lines)


def extract_pptx(path: pathlib.Path) -> str:
    prs = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f'--- Slide {i} ---')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = ''.join(run.text for run in para.runs).strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append(' | '.join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f'[notes] {notes}')
    return '\n'.join(lines)


def extract_pdf(path: pathlib.Path) -> str:
    try:
        r = subprocess.run(
            ['pdftotext', '-layout', str(path), '-'],
            capture_output=True, text=True, timeout=60,
        )
        return r.stdout
    except FileNotFoundError:
        return '[pdftotext not available]'


def extract(path: pathlib.Path) -> str:
    suffix = path.suffix.lower()
    if suffix == '.docx':
        return extract_docx(path)
    if suffix == '.pptx':
        return extract_pptx(path)
    if suffix == '.pdf':
        return extract_pdf(path)
    raise ValueError(f'Unsupported file type: {path}')


SUBFOLDERS = ['PS369_Handouts', 'PS369_LectureMaterials', 'PS369_StudyGuides']
STANDALONE_FILES = ['PS369_PostSovietPolitics_Spring2026_updated.docx']


def run() -> None:
    total_processed = 0
    total_skipped = 0

    for sub in SUBFOLDERS:
        folder = ROOT / sub
        if not folder.exists():
            print(f'Skipping missing folder: {sub}')
            continue
        out_dir = folder / 'plain_text'
        out_dir.mkdir(exist_ok=True)
        print(f'\n== {sub} ==')
        for src in sorted(folder.iterdir()):
            if src.is_dir():
                continue
            if src.suffix.lower() not in {'.docx', '.pptx', '.pdf'}:
                continue
            out_path = out_dir / (src.stem + '.txt')
            try:
                text = extract(src)
                out_path.write_text(text)
                total_processed += 1
                print(f'  {src.name} -> {len(text):,} chars')
            except Exception as e:
                total_skipped += 1
                print(f'  [skip] {src.name}: {e}')

    for fname in STANDALONE_FILES:
        src = ROOT / fname
        if not src.exists():
            continue
        out_path = ROOT / (src.stem + '_plain_text.txt')
        try:
            text = extract(src)
            out_path.write_text(text)
            total_processed += 1
            print(f'\n{src.name} -> {out_path.name} ({len(text):,} chars)')
        except Exception as e:
            total_skipped += 1
            print(f'\n[skip] {src.name}: {e}')

    print(f'\nTotal: processed {total_processed}, skipped {total_skipped}')


if __name__ == '__main__':
    run()
